from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

TEMPLATE_PATH = '/Users/aakash/Documents/Antigravity Projects/uofu-symposium/AIinBusiness_Symposium_PP Final Template.pptx'
OUTPUT_PATH = '/Users/aakash/Documents/Antigravity Projects/uofu-symposium/AI_Readiness_Symposium_Presentation.pptx'

def create_presentation():
    if not os.path.exists(TEMPLATE_PATH):
        print(f"Error: Template not found at {TEMPLATE_PATH}")
        return

    # Load template
    prs = Presentation(TEMPLATE_PATH)

    # --- Slide 1: Title ---
    # Layout 0: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = "AI in Business: Impact and Opportunity"
    
    # Subtitle (using subtitle placeholder if available, otherwise creating textbox)
    # The analysis showed Layout 0 usually has Title (0) and Subtitle (1)
    if len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = "Building Organizational Readiness for AI\n\nAakash Tuli\nDirector of Architecture"
    else:
        # Fallback if placeholder missing
        left = Inches(1)
        top = Inches(4)
        width = Inches(8)
        height = Inches(2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "Building Organizational Readiness for AI\n\nAakash Tuli\nDirector of Architecture"


    # --- Slide 2: The Strategic Pivot ---
    # Layout 1: Title and Content (Standard)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "The Strategic Pivot: Operational Alpha"
    
    # Content
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "The Scale Gap:"
    p = tf.add_paragraph()
    p.text = "80% of AI pilots never reach production."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Goal:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Move from 'AI Curiosity' to 'Operational Alpha'."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Result:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Enterprise Engine & AI ROI."
    p.level = 1
    
    # Interaction/Poll
    left = Inches(1)
    top = Inches(5.5)
    width = Inches(8)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "POLL: Can you name one AI initiative that is a permanent part of your P&L today?"
    p.font.bold = True
    # p.font.color.rgb = None # Default theme color


    # --- Slide 3: The 10-20-70 Rule ---
    # Layout 3: Two Content (Comparison)
    slide_layout = prs.slide_layouts[3] 
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "The 10-20-70 Rule of AI Success"
    
    # Left Content
    left_body = slide.placeholders[1]
    tf = left_body.text_frame
    tf.text = "Investment Strategy:"
    p = tf.add_paragraph()
    p.text = "10% Algorithms (Commodity)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "20% Technology (Infrastructure)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "70% People & Process (Transformation)"
    p.level = 1

    # Right Content
    right_body = slide.placeholders[2]
    tf = right_body.text_frame
    tf.text = "\"The model is a commodity — the transformation is human.\""
    
    # Interaction
    left = Inches(1)
    top = Inches(6)
    width = Inches(8.5)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "CHALLENGE: If your budget is 80% tech, you are building a pilot. Where is your weight?"
    p.font.bold = True


    # --- Slide 4: The Readiness Equation ---
    # Layout 1: Title and Content
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "The Readiness Equation"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Readiness = (Data × Strategy)^Culture / Risk"
    
    p = tf.add_paragraph()
    p.text = "Data: The Fuel (Clean, connected)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Strategy: The Direction"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Culture: The Multiplier (Exponential effect)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Risk: The Friction (Denominator)"
    p.level = 1

    # Interaction
    left = Inches(1)
    top = Inches(5.5)
    width = Inches(8)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "DISCUSSION: Which variable is your current 'Denominator'? Data? Strategy? Or Culture?"
    p.font.bold = True


    # --- Slide 5: Infrastructure (Silos to Fabric) ---
    # Layout 4: Comparison (Two Content with Headers typically, or Layout 3)
    # Using Layout 3 (Two Content) for side-by-side
    slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Infrastructure: From Silos to Fabric"
    
    # Left: Legacy
    left_body = slide.placeholders[1]
    tf = left_body.text_frame
    tf.text = "Legacy (Then)"
    p = tf.add_paragraph()
    p.text = "Disconnected Silos (CRM, ERP)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Batch updates"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "High latency"
    p.level = 1

    # Right: AI Fabric
    right_body = slide.placeholders[2]
    tf = right_body.text_frame
    tf.text = "AI Fabric (Now)"
    p = tf.add_paragraph()
    p.text = "Cloud Elasticity + Real-Time Streams"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Model-Agnostic Orchestration"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "\"Context in Motion\""
    p.level = 1
    
    # Interaction
    left = Inches(1)
    top = Inches(6)
    width = Inches(8.5)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "CHALLENGE: If disrupted today, could your stack respond in 48 hours?"
    p.font.bold = True


    # --- Slide 6: The Human Engine ---
    # Layout 1: Title and Content 
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "The Human Engine & Org Design"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "New Roles:"
    p = tf.add_paragraph()
    p.text = "Context Engineers (Truth Architects)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "AI Product Managers (ROI Owners)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Reskilled SMEs"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Structure (Hub & Spoke):"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "AI CoE (Standards) + Business Units (P&L)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Feedback Loops accelerate innovation"
    p.level = 1

    # Interaction
    left = Inches(1)
    top = Inches(5.5)
    width = Inches(8)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "SCALE 1-10: How confident are you in vetting AI ROI without vendor slides?"
    p.font.bold = True


    # Save
    prs.save(OUTPUT_PATH)
    print(f"Successfully saved presentation to: {OUTPUT_PATH}")

if __name__ == "__main__":
    create_presentation()
